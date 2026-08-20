import os
import sys
import yaml
import json
import re
import gc
import tempfile
import io

try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from .chunking import SLMChunker
except ImportError:
    try:
        from chunking import SLMChunker  # type: ignore
    except ImportError:
        SLMChunker = None

try:
    import onnxruntime_genai as og
except ImportError:
    og = None

_shared_ocr_engine = None

def get_ocr_engine():
    global _shared_ocr_engine
    if _shared_ocr_engine is None:
        try:
            from rapidocr_onnxruntime import RapidOCR
            _shared_ocr_engine = RapidOCR()
        except Exception:
            _shared_ocr_engine = False
    return _shared_ocr_engine if _shared_ocr_engine is not False else None

def load_config() -> tuple[dict, str]:
    config_paths = [
        os.environ.get("SLM_DOCUMENT_PARSER_CONFIG"),
        "./config.yaml",
        "../config.yaml",
        os.path.join(os.path.dirname(os.path.abspath(__file__)), "config.yaml"),
        os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "config.yaml")
    ]
    for path in config_paths:
        if path and os.path.exists(path):
            try:
                with open(path, "r") as f:
                    return yaml.safe_load(f) or {}, os.path.abspath(path)
            except Exception:
                pass
    return {}, ""

class SLMDocumentParser:
    """
    A local CPU-optimized Document Parser agent powered by lightweight ONNX runtime
    and Phi-3.5 / Qwen models. Extracts data from PDF, DOCX, PPTX, Excel, CSV, Images,
    and Markdown files into structured JSON schemas and RAG semantic graph chunks.
    """
    def __init__(self, model_path=None, cache_dir=None, n_ctx=None, n_threads=None):
        n_threads = n_threads or int(os.environ.get("SLM_DOCUMENT_PARSER_N_THREADS", os.environ.get("SLM_N_THREADS", 2)))
        self.n_ctx = n_ctx or int(os.environ.get("SLM_DOCUMENT_PARSER_N_CTX", 4096))
        cache_dir = cache_dir or os.environ.get("SLM_DOCUMENT_PARSER_CACHE_DIR")

        os.environ["OMP_NUM_THREADS"] = str(n_threads)
        os.environ["MKL_NUM_THREADS"] = str(n_threads)

        self.model = None
        self.tokenizer = None
        self._vision_parser = None
        self._chunker_inst = None

        if og is not None:
            try:
                resolved_path = self._resolve_model_path(model_path, cache_dir)
                if resolved_path and os.path.exists(resolved_path):
                    self.model_path = resolved_path
                    print(f"[SLMDocumentParser] Loading ONNX model from: {self.model_path} (threads={n_threads})...")
                    self.model = og.Model(self.model_path)
                    self.tokenizer = og.Tokenizer(self.model)
            except Exception as e:
                print(f"[SLMDocumentParser] Note: ONNX model load deferred ({e}). Operating in low-RAM fallback mode.")

    @property
    def chunker(self):
        if not hasattr(self, "_chunker_inst") or self._chunker_inst is None:
            model = getattr(self, "model", None)
            tokenizer = getattr(self, "tokenizer", None)
            self._chunker_inst = SLMChunker(model, tokenizer) if (SLMChunker and model and tokenizer) else None
        return self._chunker_inst

    def _resolve_model_path(self, model_path=None, cache_dir=None) -> str:
        if model_path:
            if not os.path.exists(model_path):
                return ""
            return os.path.abspath(model_path)

        config, config_file_path = load_config()
        model_config = config.get("models", {}).get("document_parser", {})
        config_path = model_config.get("path", "../../models/phi-3.5-mini-instruct-onnx")
        config_path = os.path.expanduser(config_path)

        if not os.path.isabs(config_path) and config_file_path:
            config_path = os.path.abspath(os.path.join(os.path.dirname(config_file_path), config_path))

        if os.path.exists(config_path):
            for root, dirs, files in os.walk(config_path):
                if "genai_config.json" in files:
                    return root
            return config_path

        return ""

    def parse_document(self, file_path: str) -> dict:
        """
        High-level document parsing API for PDF, Word, PowerPoint, Excel, Images, and Text files.
        Returns a dict containing markdown text, character count, and status.
        """
        extracted_text = self.extract_text(file_path)
        return {
            "status": "200 OK",
            "markdown": extracted_text,
            "file_path": file_path,
            "character_count": len(extracted_text)
        }

    def extract_text(self, file_path: str) -> str:
        """
        Extracts structured text from ANY document format (PDF, DOCX, PPTX, XLSX, CSV, Images, TXT).
        Uses PyMuPDF / fitz, pdfplumber, pypdf, docx, pptx, openpyxl, and RapidOCR for low RAM footprint.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Document file not found: {file_path}")

        ext = os.path.splitext(file_path)[1].lower()

        # 1. PDF Files (.pdf)
        if ext == ".pdf":
            return self._extract_pdf_text(file_path)

        # 2. Word Documents (.docx, .doc)
        if ext == ".docx":
            if Document is not None:
                try:
                    doc = Document(file_path)
                    text_blocks = []
                    for p in doc.paragraphs:
                        if p.text.strip():
                            text_blocks.append(p.text.strip())
                    for table in doc.tables:
                        for row in table.rows:
                            row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                            if row_text:
                                text_blocks.append(row_text)
                    if text_blocks:
                        return "\n".join(text_blocks)
                except Exception as e:
                    print(f"[SLMDocumentParser] Native DOCX parser error: {e}. Using zip fallback.")
            return self._extract_docx_fallback(file_path)

        if ext == ".doc":
            return self._extract_ole_text(file_path, "WordDocument")

        # 3. PowerPoint Presentations (.pptx, .ppt)
        if ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                slides_text = []
                for idx, slide in enumerate(prs.slides):
                    slide_lines = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_lines.append(shape.text.strip())
                    if slide_lines:
                        slides_text.append(f"--- Slide {idx+1} ---\n" + "\n".join(slide_lines))
                if slides_text:
                    return "\n\n".join(slides_text)
            except Exception as e:
                print(f"[SLMDocumentParser] Native PPTX parser error: {e}. Using zip fallback.")
            fallback_slides = self._extract_pptx_fallback(file_path)
            return "\n\n".join([f"--- Slide {num} ---\n{text}" for num, text in fallback_slides])

        if ext == ".ppt":
            return self._extract_ole_text(file_path, "PowerPoint Document")

        # 4. Excel & CSV Spreadsheets (.xlsx, .xls, .csv, .tsv)
        if ext in [".xlsx", ".xls"]:
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, data_only=True)
                sheet_texts = []
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    rows = []
                    for row in sheet.iter_rows(values_only=True):
                        row_vals = [str(val).strip() for val in row if val is not None and str(val).strip()]
                        if row_vals:
                            rows.append(" | ".join(row_vals))
                    if rows:
                        sheet_texts.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows[:200]))
                if sheet_texts:
                    return "\n\n".join(sheet_texts)
            except Exception as e:
                print(f"[SLMDocumentParser] Excel extraction failed: {e}")

        if ext in [".csv", ".tsv"]:
            try:
                import csv
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    reader = csv.reader(f, delimiter="\t" if ext == ".tsv" else ",")
                    rows = [" | ".join(r) for r in reader if r]
                    if rows:
                        return f"--- CSV Data ({os.path.basename(file_path)}) ---\n" + "\n".join(rows[:250])
            except Exception as e:
                print(f"[SLMDocumentParser] CSV extraction failed: {e}")

        # 5. Image Files (.png, .jpg, .jpeg, .webp, .tiff, .bmp)
        if ext in [".png", ".jpg", ".jpeg", ".webp", ".tiff", ".bmp"]:
            ocr = get_ocr_engine()
            if ocr is not None:
                try:
                    ocr_res, _ = ocr(file_path)
                    if ocr_res:
                        extracted = [item[1] for item in ocr_res if len(item) >= 2]
                        if extracted:
                            return f"--- OCR Image Text ({os.path.basename(file_path)}) ---\n" + "\n".join(extracted)
                except Exception as e:
                    print(f"[SLMDocumentParser] RapidOCR image parsing failed: {e}")

        # 6. Plain Text / Code / Markdown Fallback
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception as e:
            print(f"[SLMDocumentParser] Text fallback error: {e}")
            return ""

    def _extract_pdf_text(self, pdf_path: str) -> str:
        """Extracts PDF text using PyMuPDF (fitz), pdfplumber, pypdf, and RapidOCR ONNX fallback."""
        pages_text = []

        # Tier 1: PyMuPDF (fitz) + RapidOCR for scanned/sparse pages
        try:
            import pymupdf as fitz
            doc = fitz.open(pdf_path)
            ocr = get_ocr_engine()
            for i in range(len(doc)):
                page = doc[i]
                text = page.get_text("text").strip()
                if len(text) >= 20:
                    pages_text.append(f"--- Page {i+1} ---\n{text}")
                else:
                    ocr_text = ""
                    if ocr is not None:
                        try:
                            pix = page.get_pixmap(dpi=150)
                            img_bytes = pix.tobytes("png")
                            ocr_res, _ = ocr(img_bytes)
                            if ocr_res:
                                ocr_text = "\n".join([item[1] for item in ocr_res if len(item) >= 2]).strip()
                        except Exception:
                            pass
                    combined = (text + "\n" + ocr_text).strip() if (text and ocr_text) else (ocr_text or text)
                    if combined:
                        pages_text.append(f"--- Page {i+1} ---\n{combined}")
            doc.close()
        except Exception as e:
            print(f"[SLMDocumentParser] PyMuPDF PDF extraction failed: {e}")

        # Tier 2: pdfplumber fallback
        if not pages_text:
            try:
                import pdfplumber
                with pdfplumber.open(pdf_path) as pdf:
                    for i, page in enumerate(pdf.pages):
                        txt = page.extract_text() or ""
                        if txt.strip():
                            pages_text.append(f"--- Page {i+1} ---\n{txt.strip()}")
            except Exception as e:
                print(f"[SLMDocumentParser] pdfplumber PDF extraction failed: {e}")

        # Tier 3: pypdf fallback
        if not pages_text and PdfReader is not None:
            try:
                reader = PdfReader(pdf_path)
                for i, page in enumerate(reader.pages):
                    txt = page.extract_text() or ""
                    if txt.strip():
                        pages_text.append(f"--- Page {i+1} ---\n{txt.strip()}")
            except Exception as e:
                print(f"[SLMDocumentParser] pypdf PDF extraction failed: {e}")

        return "\n\n".join(pages_text) if pages_text else ""

    def _get_vision_parser(self):
        if not hasattr(self, "_vision_parser") or self._vision_parser is None:
            try:
                from slm_vision_parser.vision_parser import SLMVisionParser  # type: ignore
                self._vision_parser = SLMVisionParser()
            except Exception:
                self._vision_parser = None
        return self._vision_parser

    def _convert_page_to_image(self, pdf_path: str, page_idx: int) -> str:
        """Converts a specific PDF page index to a PNG image using pypdfium2."""
        try:
            import pypdfium2 as pdfium
            doc = pdfium.PdfDocument(pdf_path)
            if page_idx < 0 or page_idx >= len(doc):
                return ""
            page = doc[page_idx]
            bitmap = page.render(scale=1.2)
            pil_image = bitmap.to_pil()
            temp_dir = tempfile.gettempdir()
            img_path = os.path.join(temp_dir, f"page_{page_idx+1}_{os.path.basename(pdf_path)}.png")
            pil_image.save(img_path)
            return img_path
        except Exception as e:
            print(f"[SLMDocumentParser] Page image rendering failed: {e}")
            return ""

    def _convert_office_to_pdf(self, file_path: str) -> str:
        """Converts doc/docx/ppt/pptx to pdf via LibreOffice CLI if available."""
        import subprocess
        temp_dir = tempfile.gettempdir()
        try:
            soffice_cmd = "soffice"
            if sys.platform == "darwin":
                mac_soffice = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
                if os.path.exists(mac_soffice):
                    soffice_cmd = mac_soffice

            cmd = [soffice_cmd, "--headless", "--convert-to", "pdf", "--outdir", temp_dir, file_path]
            subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True, timeout=30)
            basename = os.path.splitext(os.path.basename(file_path))[0] + ".pdf"
            pdf_path = os.path.join(temp_dir, basename)
            if os.path.exists(pdf_path):
                return pdf_path
        except Exception as e:
            print(f"[SLMDocumentParser] LibreOffice conversion note: {e}")
        return ""

    def _parse_page_visually(self, image_path: str) -> dict:
        """Parses a page image using RapidOCR or Vision Parser."""
        if not image_path or not os.path.exists(image_path):
            return {"page_text": "", "tables": [], "captions": []}

        ocr = get_ocr_engine()
        ocr_text = ""
        if ocr is not None:
            try:
                ocr_res, _ = ocr(image_path)
                if ocr_res:
                    ocr_text = "\n".join([item[1] for item in ocr_res if len(item) >= 2])
            except Exception as e:
                print(f"[SLMDocumentParser] Visual page OCR failed: {e}")

        return {
            "page_text": ocr_text,
            "tables": [],
            "captions": []
        }

    def _assemble_and_correct_page_markdown(self, page_data: dict, page_num: int) -> str:
        """Uses ONNX model or structured formatting to assemble raw page OCR text into clean Markdown."""
        page_text = page_data.get("page_text", "") or ""
        tables_str = "\n\n".join([f"Table Area Description:\n{t}" for t in page_data.get("tables", []) if t])
        captions_str = "\n\n".join([f"Image/Chart Description: {c}" for c in page_data.get("captions", []) if c])

        if not self.model or not self.tokenizer or og is None:
            parts = [f"--- Page {page_num} ---", page_text]
            if tables_str:
                parts.append(tables_str)
            if captions_str:
                parts.append(captions_str)
            return "\n\n".join([p for p in parts if p.strip()]).strip()

        system_prompt = (
            "You are a local Document Structuring Agent.\n"
            "Format the raw page OCR text into clean structured Markdown.\n"
            "Maintain correct heading hierarchy (use #, ##, ###). Do not include preamble or chatter, output ONLY Markdown."
        )

        user_prompt = (
            f"--- Page {page_num} ---\n"
            f"Raw OCR Text:\n{page_text}\n\n"
            f"{tables_str}\n\n"
            f"{captions_str}\n\n"
            "Output structured Markdown:"
        )

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]

        full_prompt = ""
        for msg in messages:
            full_prompt += f"<|{msg['role']}|>\n{msg['content']}<|end|>\n"
        full_prompt += "<|assistant|>\n"

        try:
            input_tokens = self.tokenizer.encode(full_prompt)
            params = og.GeneratorParams(self.model)
            params.set_search_options(max_length=len(input_tokens) + 512, temperature=0.7)

            generator = og.Generator(self.model, params)
            generator.append_tokens(input_tokens)

            response_text = ""
            while not generator.is_done():
                generator.generate_next_token()
                new_tokens = generator.get_next_tokens()
                if len(new_tokens) > 0:
                    token_id = int(new_tokens[0])
                    if token_id in (151643, 151645, 248046, 248044, 248045, 32000, 32007):
                        break
                    response_text += self.tokenizer.decode(new_tokens)
            return response_text.strip()
        except Exception as e:
            print(f"[SLMDocumentParser] LLM page markdown correction error: {e}")
            return page_text.strip()

    def _segment_and_extract_chunks(self, full_markdown: str, source_name: str) -> list[dict]:
        """Uses Phi-3.5 ONNX or fallback chunker to segment markdown text into structured semantic chunks."""
        if self.chunker:
            try:
                return self.chunker.segment_and_extract_chunks(full_markdown, source_name)
            except Exception as e:
                print(f"[SLMDocumentParser] Chunker segment error: {e}. Falling back.")
        return self._fallback_semantic_chunker(full_markdown, source_name)

    def _fallback_semantic_chunker(self, text: str, source_name: str) -> list[dict]:
        """Fallback chunker that splits text by paragraphs and extracts basic headers."""
        if self.chunker:
            return self.chunker.fallback_semantic_chunker(text, source_name)
        paragraphs = text.split("\n\n")
        chunks = []
        current_h1 = ""
        current_h2 = ""

        for idx, para in enumerate(paragraphs):
            para = para.strip()
            if not para:
                continue

            if para.startswith("# "):
                current_h1 = para[2:].strip()
                continue
            elif para.startswith("## "):
                current_h2 = para[3:].strip()
                continue
            elif para.startswith("### "):
                current_h2 = para[4:].strip()
                continue

            product = ""
            words = para.split()
            for w in words:
                w_clean = re.sub(r'[^\w]', '', w)
                if w_clean.istitle() and len(w_clean) > 4:
                    product = w_clean
                    break

            if len(para.split()) >= 3:
                chunks.append({
                    "text": para,
                    "metadata": {
                        "source": os.path.basename(source_name),
                        "heading": current_h1,
                        "subheading": current_h2,
                        "product": product,
                        "key_terms": [current_h1] if current_h1 else ["general"],
                        "format": os.path.splitext(source_name)[1][1:].lower()
                    }
                })
        return chunks

    def _link_chunks(self, chunks: list[dict]) -> list[dict]:
        """Post-processes chunks to link related section siblings and product references."""
        if self.chunker:
            return self.chunker.link_chunks(chunks)

        for idx, chunk in enumerate(chunks):
            chunk["metadata"]["chunk_index"] = idx
            chunk["metadata"]["related_chunks"] = []

        for i, c1 in enumerate(chunks):
            meta1 = c1["metadata"]
            idx1 = meta1["chunk_index"]
            related = set()

            for j in range(len(chunks)):
                if i == j:
                    continue
                c2 = chunks[j]
                meta2 = c2["metadata"]
                idx2 = meta2["chunk_index"]

                if meta1.get("heading") and meta1.get("heading") == meta2.get("heading"):
                    related.add(idx2)
                if meta1.get("subheading") and meta1.get("subheading") == meta2.get("subheading"):
                    related.add(idx2)
                if meta1.get("product") and meta1.get("product").lower() == meta2.get("product", "").lower():
                    related.add(idx2)

                kt1 = set(meta1.get("key_terms", []))
                kt2 = set(meta2.get("key_terms", []))
                shared_terms = kt1.intersection(kt2)
                if len(shared_terms) >= 2:
                    related.add(idx2)

            meta1["related_chunks"] = sorted(list(related))
        return chunks

    def parse_and_chunk_stream(self, file_path: str, chunk_size: int = 500, chunk_overlap: int = 50):
        """Streaming parser that yields linked semantic chunks page-by-page."""
        full_text = self.extract_text(file_path)
        semantic_chunks = self._segment_and_extract_chunks(full_text, file_path)
        self._link_chunks(semantic_chunks)
        for chunk in semantic_chunks:
            yield chunk

    def chunk_document(self, file_path: str, chunk_size: int = 500, chunk_overlap: int = 50) -> list[dict]:
        """Runs the streaming parser and returns the complete list of linked semantic chunks."""
        return list(self.parse_and_chunk_stream(file_path, chunk_size, chunk_overlap))

    def export_chunks_to_excel(self, chunks: list[dict], output_path: str, append: bool = False) -> None:
        """Exports a list of chunk dicts to an Excel (.xlsx) file using openpyxl."""
        try:
            import openpyxl
        except ImportError:
            raise ImportError("openpyxl is required to export to Excel. Please run: pip install openpyxl")

        if append and os.path.exists(output_path):
            try:
                wb = openpyxl.load_workbook(output_path)
                ws = wb.active
                max_idx = -1
                for r in range(2, ws.max_row + 1):
                    val = ws.cell(row=r, column=1).value
                    if val is not None:
                        try:
                            max_idx = max(max_idx, int(val))
                        except ValueError:
                            pass
                start_idx = max_idx + 1
            except Exception as e:
                print(f"[SLMDocumentParser] Failed to load existing Excel file: {e}. Creating new sheet.")
                wb = openpyxl.Workbook()
                ws = wb.active
                ws.title = "RAG Chunks"
                ws.append(["Chunk Index", "Source File", "Heading", "Subheading", "Product", "Related Chunks", "Text"])
                start_idx = 0
        else:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "RAG Chunks"
            ws.append(["Chunk Index", "Source File", "Heading", "Subheading", "Product", "Related Chunks", "Text"])
            start_idx = 0

        current_idx = start_idx
        for chunk in chunks:
            metadata = chunk.get("metadata", {})
            related_str = ",".join(map(str, metadata.get("related_chunks", [])))
            row = [
                current_idx,
                metadata.get("source", ""),
                metadata.get("heading", ""),
                metadata.get("subheading", ""),
                metadata.get("product", ""),
                related_str,
                chunk.get("text", "")
            ]
            ws.append(row)
            current_idx += 1

        wb.save(output_path)
        print(f"[SLMDocumentParser] Exported {len(chunks)} chunks to Excel at: {output_path}")

    def _extract_docx_fallback(self, file_path: str) -> str:
        import zipfile
        try:
            with zipfile.ZipFile(file_path, 'r') as z:
                content = z.read('word/document.xml').decode('utf-8', errors='ignore')
                text_segments = re.findall(r'<w:t[^>]*>(.*?)</w:t>', content)
                text = "\n".join(text_segments)
                for k, v in [("&amp;", "&"), ("&lt;", "<"), ("&gt;", ">"), ("&quot;", '"'), ("&apos;", "'")]:
                    text = text.replace(k, v)
                return text
        except Exception as e:
            print(f"[SLMDocumentParser] Fallback DOCX extraction error: {e}")
            return ""

    def _extract_pptx_fallback(self, file_path: str) -> list[tuple[int, str]]:
        import zipfile
        slides = []
        try:
            with zipfile.ZipFile(file_path, 'r') as z:
                slide_files = [f for f in z.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
                def extract_slide_num(name):
                    match = re.search(r'slide(\d+)\.xml', name)
                    return int(match.group(1)) if match else 0
                slide_files.sort(key=extract_slide_num)

                for file_name in slide_files:
                    slide_num = extract_slide_num(file_name)
                    content = z.read(file_name).decode('utf-8', errors='ignore')
                    text_segments = re.findall(r'<a:t[^>]*>(.*?)</a:t>', content)
                    text = " ".join(text_segments).strip()
                    if text:
                        for k, v in [("&amp;", "&"), ("&lt;", "<"), ("&gt;", ">"), ("&quot;", '"'), ("&apos;", "'")]:
                            text = text.replace(k, v)
                        slides.append((slide_num, text))
        except Exception as e:
            print(f"[SLMDocumentParser] Fallback PPTX extraction error: {e}")
        return slides

    def _extract_ole_text(self, file_path: str, stream_name: str) -> str:
        try:
            import olefile
            if olefile.isOleFile(file_path):
                ole = olefile.OleFileIO(file_path)
                if ole.exists(stream_name):
                    stream = ole.openstream(stream_name)
                    data = stream.read()
                    return self._extract_strings_from_bytes(data)
        except Exception as e:
            print(f"[SLMDocumentParser] OLE stream extraction error: {e}")

        try:
            with open(file_path, "rb") as f:
                data = f.read()
            return self._extract_strings_from_bytes(data)
        except Exception as e:
            print(f"[SLMDocumentParser] OLE binary fallback error: {e}")
            return ""

    def _extract_strings_from_bytes(self, data: bytes, min_len: int = 4) -> str:
        ascii_pattern = re.compile(rb'[\x20-\x7E\x0A\x0D\x09]{' + str(min_len).encode() + rb',}')
        ascii_matches = ascii_pattern.findall(data)

        utf16_pattern = re.compile(rb'(?:[\x20-\x7E\x0A\x0D\x09]\x00){' + str(min_len).encode() + rb',}')
        utf16_matches = utf16_pattern.findall(data)

        decoded = []
        for m in ascii_matches:
            try:
                decoded.append(m.decode('ascii').strip())
            except Exception:
                pass
        for m in utf16_matches:
            try:
                decoded.append(m.decode('utf-16le').strip())
            except Exception:
                pass

        return "\n".join([s for s in decoded if len(s) > 3])

    def _extract_json(self, text: str) -> str:
        match = re.search(r"```json\s*(.*?)\s*```", text, re.DOTALL)
        if match:
            return match.group(1).strip()
        brace_match = re.search(r"(\{.*\})", text, re.DOTALL)
        if brace_match:
            return brace_match.group(1).strip()
        return text.strip()

    def parse(self, file_path: str, schema_dict: dict, max_retries: int = 3, system_prompt: str = None, user_input: str = None) -> dict:
        """Parses a document file into a structured JSON dict matching schema_dict."""
        raw_text = self.extract_text(file_path)

        if not self.model or not self.tokenizer or og is None:
            # Low-RAM rule-based fallback JSON extractor
            result = {}
            for k, v in schema_dict.items():
                if isinstance(v, list):
                    result[k] = []
                elif isinstance(v, dict):
                    result[k] = {}
                elif isinstance(v, int):
                    result[k] = 0
                elif isinstance(v, float):
                    result[k] = 0.0
                elif isinstance(v, bool):
                    result[k] = False
                else:
                    match = re.search(rf"{k}[:=]\s*([^\n,;]+)", raw_text, re.IGNORECASE)
                    if match:
                        result[k] = match.group(1).strip()
                    else:
                        result[k] = f"Extracted {k} from {os.path.basename(file_path)}"
            return result

        system_prompt = (
            "You are a local Document Parser agent.\n"
            "Analyze the document text and extract the data to populate a structured JSON block matching the target schema. "
            "Return the final completed JSON inside a ```json ... ``` code block."
        )

        user_prompt = (
            f"Document Text:\n{raw_text[:8000]}\n\n"
            f"Target JSON Schema Structure to Populate:\n{json.dumps(schema_dict, indent=2)}"
        )

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]

        for attempt in range(max_retries):
            full_prompt = ""
            for msg in messages:
                full_prompt += f"<|{msg['role']}|>\n{msg['content']}<|end|>\n"
            full_prompt += "<|assistant|>\n"

            try:
                input_tokens = self.tokenizer.encode(full_prompt)
                params = og.GeneratorParams(self.model)
                params.set_search_options(max_length=len(input_tokens) + 1024, temperature=0.7)

                generator = og.Generator(self.model, params)
                generator.append_tokens(input_tokens)

                response_text = ""
                while not generator.is_done():
                    generator.generate_next_token()
                    new_tokens = generator.get_next_tokens()
                    if len(new_tokens) > 0:
                        token_id = int(new_tokens[0])
                        if token_id in (151643, 151645, 248046, 248044, 248045, 32000, 32007):
                            break
                        response_text += self.tokenizer.decode(new_tokens)

                json_block = self._extract_json(response_text)
                parsed = json.loads(json_block)
                return parsed
            except Exception as e:
                messages.append({"role": "assistant", "content": response_text if 'response_text' in locals() else ""})
                messages.append({
                    "role": "user",
                    "content": f"JSON parsing failed with error: {e}. Correct the JSON format and return the updated block inside ```json ```."
                })

        return {"error": "Failed to parse document complying with schema within retries limit"}
